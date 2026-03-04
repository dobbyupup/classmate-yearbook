from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(90, 90, 90)
LIGHT_GRAY = RGBColor(180, 180, 180)
WHITE = RGBColor(255, 255, 255)


def set_text_style(run, size=24, bold=False, color=BLACK, font="Arial"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_text_style(run, size=42, bold=True)
    p.alignment = PP_ALIGN.LEFT


def add_subtitle(slide, text, top=1.45):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_text_style(run, size=20, color=GRAY)
    p.alignment = PP_ALIGN.LEFT


def add_line(slide, top=1.35):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()


def add_bullets(slide, items, top=2.0, left=1.0, width=11.0, height=4.5, size=24):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            set_text_style(run, size=size, color=BLACK)


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.3), Inches(0.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_text_style(run, size=14, color=GRAY)
    p.alignment = PP_ALIGN.LEFT


def add_two_column(slide, left_title, left_items, right_title, right_items):
    lt = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.4), Inches(0.5))
    ltf = lt.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = left_title
    set_text_style(lr, size=24, bold=True)

    rt = slide.shapes.add_textbox(Inches(6.7), Inches(1.8), Inches(5.4), Inches(0.5))
    rtf = rt.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = right_title
    set_text_style(rr, size=24, bold=True)

    add_bullets(slide, left_items, top=2.4, left=0.9, width=5.6, height=3.8, size=20)
    add_bullets(slide, right_items, top=2.4, left=6.7, width=5.4, height=3.8, size=20)


def add_center_sentence(slide, text):
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.0), Inches(1.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_text_style(run, size=36, bold=True)


def add_small_note(slide, text):
    box = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.0), Inches(0.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_text_style(run, size=16, color=GRAY)


def add_timeline(slide, steps):
    # Horizontal baseline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.6), Inches(10.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()

    x_positions = [1.1, 3.6, 6.1, 8.6]
    for i, (title, desc) in enumerate(steps):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_positions[i]), Inches(3.45), Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLACK
        circle.line.fill.background()

        box = slide.shapes.add_textbox(Inches(x_positions[i] - 0.2), Inches(2.2), Inches(2.4), Inches(1.1))
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = title
        for run in p1.runs:
            set_text_style(run, size=18, bold=True)

        p2 = tf.add_paragraph()
        p2.text = desc
        for run in p2.runs:
            set_text_style(run, size=14, color=GRAY)


def build_ppt(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Cover
    s = prs.slides.add_slide(blank)
    add_title(s, "杨浦高校NFC毕业文创礼")
    add_line(s, top=1.35)
    add_subtitle(s, "让毕业不散场，让人才有归航", top=1.6)
    add_footer(s, "杨浦区青年人才数字基础设施提案 · 2026")

    # 2 One sentence definition
    s = prs.slides.add_slide(blank)
    add_title(s, "一句话定义")
    add_center_sentence(s, "全国首个“电子同学录+校友社群+政策触达+城市记忆”NFC毕业文创礼")
    add_small_note(s, "每位毕业生一枚NFC芯片，绑定终身ID；一碰即进入个人专属数字档案与校友网络")

    # 3 Why now
    s = prs.slides.add_slide(blank)
    add_title(s, "为什么是杨浦，为什么是现在")
    add_bullets(
        s,
        [
            "1. 全国密度最高的顶尖高校集群：复旦、同济、上财等",
            "2. “知识杨浦”品牌在青年人才中已形成情感认同",
            "3. 越早启动，越快形成数据资产与社群先发优势",
        ],
        top=2.0,
        size=24,
    )
    add_footer(s, "这不是毕业礼物，而是面向未来20年的青年人才基础设施。")

    # 4 Problems
    s = prs.slides.add_slide(blank)
    add_title(s, "当前痛点")
    add_bullets(
        s,
        [
            "• 纸质同学录易丢失、不可更新",
            "• 校友关系毕业即“断联”",
            "• 政策触达效率低，仍是“人找政策”",
            "• 人才流向缺少实时、可信的数据抓手",
        ],
        top=2.1,
        size=26,
    )

    # 5 Solution overview
    s = prs.slides.add_slide(blank)
    add_title(s, "解决方案总览")
    add_bullets(
        s,
        [
            "电子同学录",
            "校友活动广场",
            "政策直达通道",
            "人才流向数据后台（脱敏）",
            "城市记忆入口",
        ],
        top=2.0,
        size=24,
    )
    add_footer(s, "情感连接 + 服务连接 + 数据连接，三位一体。")

    # 6 Module 1
    s = prs.slides.add_slide(blank)
    add_title(s, "模块① 电子同学录")
    add_subtitle(s, "从一次填写，到终身可更新", top=1.5)
    add_bullets(
        s,
        [
            "• 唯一NFC身份绑定，终身ID",
            "• 语音留言祝福、视频永久保存、毕业照共创",
            "• 联系方式/职业轨迹/家庭动态持续更新",
            "• 十年后仍可精准找回同窗关系",
        ],
        top=2.2,
        size=23,
    )

    # 7 Module 2
    s = prs.slides.add_slide(blank)
    add_title(s, "模块② 校友活动广场")
    add_subtitle(s, "打破校门围墙，形成跨校社群资产", top=1.5)
    add_bullets(
        s,
        [
            "• 跨校联动活动：行业沙龙、创业对接、城市聚会、联谊",
            "• 支持全国异地校友发起活动",
            "• “大学同学”升级为可运营社群网络",
            "• 增强“杨浦校友圈”归属感与活跃度",
        ],
        top=2.2,
        size=23,
    )

    # 8 Module 3
    s = prs.slides.add_slide(blank)
    add_title(s, "模块③ 政策直达通道")
    add_subtitle(s, "从“人找政策”到“政策找人”", top=1.5)
    add_bullets(
        s,
        [
            "• 人才政策、创业扶持、安居补贴精准推送",
            "• 依据就业/创业/回沪意向进行点对点触达",
            "• 一键直达申请入口，提高政策转化率",
            "• 支撑“毕业即入职杨浦”的转化链路",
        ],
        top=2.2,
        size=23,
    )

    # 9 Module 4+5
    s = prs.slides.add_slide(blank)
    add_title(s, "模块④+⑤ 数据后台与城市记忆")
    add_two_column(
        s,
        "人才流向数据后台",
        [
            "• 全流程脱敏聚合",
            "• 留沪率、回沪趋势、行业分布实时可视",
            "• 触达效果可追踪，支撑政策评估",
        ],
        "城市记忆入口",
        [
            "• 个性化杨浦打卡地图",
            "• 高校周边文化与历史人文内容",
            "• 每位毕业生专属“杨浦记忆页”",
        ],
    )

    # 10 Government value
    s = prs.slides.add_slide(blank)
    add_title(s, "对政府的核心价值")
    add_bullets(
        s,
        [
            "1. 实时看见人才流向，提升治理精度",
            "2. 政策直通提升转化，增强人才粘性",
            "3. 打造全国数字毕业文化IP，提升城市影响力",
            "4. 促进回流消费、创业、置业，拉动经济与文化软实力",
        ],
        top=2.0,
        size=24,
    )

    # 11 Communication
    s = prs.slides.add_slide(blank)
    add_title(s, "传播与出圈路径")
    add_bullets(
        s,
        [
            "产品发布  →  高校共创内容  →  抖音话题发酵  →  城市品牌出圈  →  校友持续回流",
            "",
            "#杨浦毕业礼太卷了   #杨浦大学城数字毕业季   #毕业不散场",
        ],
        top=2.4,
        size=22,
    )

    # 12 Roadmap
    s = prs.slides.add_slide(blank)
    add_title(s, "时间推进表（3-6月）")
    add_timeline(
        s,
        [
            ("3-4月", "方案确认、IP设计、NFC打样"),
            ("4-5月", "高校试点对接（复旦/同济）"),
            ("5-6月", "批量生产 + 抖音预热"),
            ("6月", "毕业季发放 + 话题引爆"),
        ],
    )
    add_footer(s, "后续复制到全国城市，形成杨浦IP先发优势。")

    # 13 Growth
    s = prs.slides.add_slide(blank)
    add_title(s, "增长预估")
    add_center_sentence(s, "年新增约7万节点 · 5年累计约35万活跃校友节点")
    add_small_note(s, "节点越多，政策触达效率、人才回流概率、城市品牌势能越强")

    # 14 Risk control
    s = prs.slides.add_slide(blank)
    add_title(s, "风险与保障")
    add_bullets(
        s,
        [
            "• 数据安全：脱敏处理、分级权限、合规审计",
            "• 内容治理：实名认证与社区风控机制",
            "• 运营持续：校地共建 + 市场化运营团队",
            "• 成本控制：先试点后扩容，分阶段投入",
        ],
        top=2.1,
        size=24,
    )

    # 15 Closing
    s = prs.slides.add_slide(blank)
    add_center_sentence(s, "让每一份毕业纪念，变成杨浦未来的人才连接器")
    add_small_note(s, "提案单位 / 日期 / 联系方式")

    # Global white background and minimalist line accents
    for slide in prs.slides:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # top thin line
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(0.18), Inches(11.7), Inches(0.01))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = LIGHT_GRAY
        top_line.line.fill.background()

    prs.save(output_path)


if __name__ == "__main__":
    build_ppt("杨浦高校NFC毕业文创礼_黑白极简版.pptx")
